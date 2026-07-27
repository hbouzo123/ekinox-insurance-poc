import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors
COLOR_PRIMARY_RGB = RGBColor(134, 49, 142)    # Ekinox Purple
COLOR_SECONDARY_RGB = RGBColor(13, 148, 136)  # Teal
COLOR_DARK_RGB = RGBColor(31, 41, 55)         # Charcoal
COLOR_LIGHT_RGB = RGBColor(249, 250, 251)      # Light Gray
COLOR_BORDER_RGB = RGBColor(229, 231, 235)     # Very Light Gray border
COLOR_WHITE_RGB = RGBColor(255, 255, 255)
COLOR_MUTED_RGB = RGBColor(107, 114, 128)      # Muted Gray

# Slide Dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def apply_background(slide, color_rgb):
    """Draw a full slide rectangle to set background color."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color_rgb
    bg.line.fill.background() # No border
    # Send to back by removing and re-inserting at index 0 (not easily exposed in API, so we just add it first)

def add_header(slide, title, category="EKINOX INSURANCE PLATFORM"):
    """Add category tracker and clean slide title."""
    # Category Tracker
    tx_cat = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(0.4))
    tf_cat = tx_cat.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_bottom = tf_cat.margin_right = Inches(0)
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY_RGB
    
    # Title
    tx_title = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12.0), Inches(0.7))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_bottom = tf_title.margin_right = Inches(0)
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY_RGB
    
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BORDER_RGB
    line.line.fill.background()

def add_footer(slide, current_page, total_pages=15):
    """Add footer branding and page numbering."""
    tx_foot = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.3))
    tf_foot = tx_foot.text_frame
    tf_foot.margin_left = tf_foot.margin_top = tf_foot.margin_bottom = tf_foot.margin_right = Inches(0)
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = f"Ekinox Consulting  |  Proposition de Transformation  |  Confidentialité SanlamAllianz  |  Slide {current_page}"
    p_foot.font.name = "Arial"
    p_foot.font.size = Pt(8.5)
    p_foot.font.color.rgb = COLOR_MUTED_RGB
    p_foot.alignment = PP_ALIGN.RIGHT

def create_title_slide(prs, title, subtitle, category="PROPOSITION PARTENAIRE"):
    """Generate a premium dark title slide."""
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Dark purple background
    apply_background(slide, COLOR_PRIMARY_RGB)
    
    # Category / Tag
    tx_cat = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(0.4))
    tf_cat = tx_cat.text_frame
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_SECONDARY_RGB
    
    # Title
    tx_title = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.8))
    tf_title = tx_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = "Arial"
    p_title.font.size = Pt(40)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_WHITE_RGB
    
    # Subtitle
    tx_sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.3), Inches(0.8))
    tf_sub = tx_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(16)
    p_sub.font.italic = True
    p_sub.font.color.rgb = COLOR_LIGHT_RGB
    
    # Accent strip in the center
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(5.2), Inches(3.0), Inches(0.06))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_SECONDARY_RGB
    strip.line.fill.background()
    
    # Bottom metadata
    tx_meta = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.0))
    tf_meta = tx_meta.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Destinataire : COMEX SanlamAllianz Group\nDate : Juillet 2026  |  Statut : STRICTEMENT CONFIDENTIEL"
    p_meta.font.name = "Calibri"
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = COLOR_LIGHT_RGB

def create_closing_slide(prs, category):
    """Generate a clean dark closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, COLOR_PRIMARY_RGB)
    
    # Title
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(11.333), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Merci pour votre confiance"
    p.font.name = "Arial"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE_RGB
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = f"Co-construisons le futur de l'assurance chez SanlamAllianz"
    p2.font.name = "Arial"
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_SECONDARY_RGB
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    
    # Contact metadata
    tx_meta = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(0.8))
    tf_meta = tx_meta.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = "Ekinox Consulting Team  |  contact@ekinox.io  |  Projet Pilote Accelerator"
    p_meta.font.name = "Calibri"
    p_meta.font.size = Pt(11)
    p_meta.font.color.rgb = COLOR_LIGHT_RGB
    p_meta.alignment = PP_ALIGN.CENTER

def add_split_card_slide(prs, title, left_title, left_text, right_title, right_text, category, page_num):
    """Generate a side-by-side card slide (ideal for client/internal splits)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, COLOR_WHITE_RGB)
    add_header(slide, title, category)
    add_footer(slide, page_num)
    
    # Positions
    y = Inches(1.8)
    h = Inches(4.8)
    w = Inches(5.8)
    x1 = Inches(0.6)
    x2 = Inches(6.933)
    
    # Card 1 (Left - Client)
    card1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y, w, h)
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_LIGHT_RGB
    card1.line.color.rgb = COLOR_BORDER_RGB
    
    # Simulated thick left border (Secondary color)
    border1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y, Inches(0.08), h)
    border1.fill.solid()
    border1.fill.fore_color.rgb = COLOR_PRIMARY_RGB
    border1.line.fill.background()
    
    tx1 = slide.shapes.add_textbox(x1 + Inches(0.2), y + Inches(0.2), w - Inches(0.4), h - Inches(0.4))
    tf1 = tx1.text_frame
    tf1.word_wrap = True
    p1_title = tf1.paragraphs[0]
    p1_title.text = left_title.upper()
    p1_title.font.name = "Arial"
    p1_title.font.size = Pt(14)
    p1_title.font.bold = True
    p1_title.font.color.rgb = COLOR_PRIMARY_RGB
    p1_title.space_after = Pt(14)
    
    # Ingest text lines or paragraph splits
    lines1 = left_text.split("\n\n")
    for line in lines1:
        if line.strip():
            p = tf1.add_paragraph()
            p.text = line.strip()
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_DARK_RGB
            p.space_after = Pt(8)
            p.line_spacing = 1.15
            
    # Card 2 (Right - Internal)
    card2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2, y, w, h)
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_LIGHT_RGB
    card2.line.color.rgb = COLOR_BORDER_RGB
    
    # Simulated thick left border (Primary color)
    border2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2, y, Inches(0.08), h)
    border2.fill.solid()
    border2.fill.fore_color.rgb = COLOR_SECONDARY_RGB
    border2.line.fill.background()
    
    tx2 = slide.shapes.add_textbox(x2 + Inches(0.2), y + Inches(0.2), w - Inches(0.4), h - Inches(0.4))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2_title = tf2.paragraphs[0]
    p2_title.text = right_title.upper()
    p2_title.font.name = "Arial"
    p2_title.font.size = Pt(14)
    p2_title.font.bold = True
    p2_title.font.color.rgb = COLOR_SECONDARY_RGB
    p2_title.space_after = Pt(14)
    
    lines2 = right_text.split("\n\n")
    for line in lines2:
        if line.strip():
            p = tf2.add_paragraph()
            p.text = line.strip()
            p.font.name = "Calibri"
            p.font.size = Pt(10.5)
            p.font.color.rgb = COLOR_DARK_RGB
            p.space_after = Pt(6)
            p.line_spacing = 1.1

def add_grid_slide(prs, title, intro_text, items, category, page_num):
    """Generate a slide with a 3x2 grid of cards (ideal for Product Blueprint)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, COLOR_WHITE_RGB)
    add_header(slide, title, category)
    add_footer(slide, page_num)
    
    # Intro sentence
    if intro_text:
        tx_intro = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.5))
        tf_intro = tx_intro.text_frame
        tf_intro.margin_left = tf_intro.margin_top = tf_intro.margin_bottom = tf_intro.margin_right = Inches(0)
        p_intro = tf_intro.paragraphs[0]
        p_intro.text = intro_text
        p_intro.font.name = "Calibri"
        p_intro.font.size = Pt(11)
        p_intro.font.italic = True
        p_intro.font.color.rgb = COLOR_DARK_RGB
        y_start = Inches(2.1)
        h_card = Inches(2.2)
    else:
        y_start = Inches(1.8)
        h_card = Inches(2.4)
        
    # Grid math (3 cols x 2 rows)
    w_card = Inches(3.85)
    gap_x = Inches(0.29)
    gap_y = Inches(0.25)
    
    for idx, item in enumerate(items[:6]): # Limit to 6
        col = idx % 3
        row = idx // 3
        
        x = Inches(0.6) + col * (w_card + gap_x)
        y = y_start + row * (h_card + gap_y)
        
        # Draw background
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w_card, h_card)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT_RGB
        card.line.color.rgb = COLOR_BORDER_RGB
        
        # Left border line
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), h_card)
        border.fill.solid()
        border.fill.fore_color.rgb = COLOR_PRIMARY_RGB if idx % 2 == 0 else COLOR_SECONDARY_RGB
        border.line.fill.background()
        
        tx_card = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), w_card - Inches(0.25), h_card - Inches(0.3))
        tf_card = tx_card.text_frame
        tf_card.word_wrap = True
        p_name = tf_card.paragraphs[0]
        p_name.text = item["name"].upper()
        p_name.font.name = "Arial"
        p_name.font.size = Pt(11.5)
        p_name.font.bold = True
        p_name.font.color.rgb = COLOR_PRIMARY_RGB if idx % 2 == 0 else COLOR_SECONDARY_RGB
        p_name.space_after = Pt(8)
        
        p_desc = tf_card.add_paragraph()
        p_desc.text = item["description"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = COLOR_DARK_RGB
        p_desc.line_spacing = 1.1

def add_table_slide(prs, title, headers, data, col_widths, category, page_num):
    """Generate a slide dedicated to structured table content (Roadmap or ROI)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, COLOR_WHITE_RGB)
    add_header(slide, title, category)
    add_footer(slide, page_num)
    
    rows = len(data) + 1
    cols = len(headers)
    
    left = Inches(0.6)
    top = Inches(1.8)
    width = Inches(12.133)
    height = Inches(4.5)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set custom column widths
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = Inches(w)
            
    # Style Header
    for c_idx, head in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = head
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY_RGB
        
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.name = "Arial"
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE_RGB
        
    # Style Data
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT_RGB if r_idx % 2 == 1 else COLOR_WHITE_RGB
            
            p = cell.text_frame.paragraphs[0]
            if "€" in str(val) or "%" in str(val) or str(val).replace(" ", "").isdigit():
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT
                
            if len(p.runs) > 0:
                run = p.runs[0]
                run.font.name = "Calibri"
                run.font.size = Pt(10.5)
                run.font.color.rgb = COLOR_DARK_RGB

def add_context_slide(prs, title, overview_text, points, category, page_num):
    """Generate a clean bullet/context slide for general messaging."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide, COLOR_WHITE_RGB)
    add_header(slide, title, category)
    add_footer(slide, page_num)
    
    # Left description block
    tx_left = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.5))
    tf_left = tx_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_top = tf_left.margin_bottom = tf_left.margin_right = Inches(0)
    
    # Clean text paragraph splitting
    paragraphs = overview_text.split("\n\n")
    for idx, p_text in enumerate(paragraphs):
        if p_text.strip():
            p = tf_left.paragraphs[0] if idx == 0 else tf_left.add_paragraph()
            p.text = p_text.strip()
            p.font.name = "Arial"
            p.font.size = Pt(13)
            p.font.color.rgb = COLOR_DARK_RGB
            p.space_after = Pt(14)
            p.line_spacing = 1.2
            
    # Right cards block
    y_start = Inches(1.8)
    h_card = Inches(1.3)
    gap_y = Inches(0.3)
    
    for idx, pt in enumerate(points[:3]): # Max 3 cards
        y = y_start + idx * (h_card + gap_y)
        
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.933), y, Inches(5.8), h_card)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT_RGB
        card.line.color.rgb = COLOR_BORDER_RGB
        
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.933), y, Inches(0.07), h_card)
        border.fill.solid()
        border.fill.fore_color.rgb = COLOR_SECONDARY_RGB
        border.line.fill.background()
        
        tx_card = slide.shapes.add_textbox(Inches(7.15), y + Inches(0.15), Inches(5.4), h_card - Inches(0.3))
        tf_card = tx_card.text_frame
        tf_card.word_wrap = True
        
        p_title = tf_card.paragraphs[0]
        p_title.text = pt["title"].upper()
        p_title.font.name = "Arial"
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY_RGB
        p_title.space_after = Pt(4)
        
        p_desc = tf_card.add_paragraph()
        p_desc.text = pt["description"]
        p_desc.font.name = "Calibri"
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = COLOR_DARK_RGB
        p_desc.line_spacing = 1.1

def build_sales_presentation(data):
    """Compile the PPTX deck for the Digital Sales Accelerator."""
    print("Building PPTX for Ekinox Digital Sales Accelerator...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    category = "DIGITAL SALES ACCELERATOR"
    
    # Slide 1: Cover
    create_title_slide(
        prs=prs,
        title=data["title"],
        subtitle="Transformer l'acquisition commerciale de SanlamAllianz via une plateforme autonome",
        category="Ekinox Insurance Platform"
    )
    
    # Slide 2: Context & Positioning
    overview = (
        "SanlamAllianz doit faire évoluer sa capacité d'acquisition client.\n\n"
        "Plutôt que d'intégrer des chatbots isolés ou d'engager d'importants budgets de refonte CRM, "
        "le Digital Sales Accelerator permet de déployer un parcours omnicanal performant, "
        "découplé du Core Insurance pour une activation immédiate."
    )
    points = [
        {"title": "Productivité commerciale", "description": "Qualification automatique des dossiers prospects avant transmission aux conseillers."},
        {"title": "Expérience Omnicanale", "description": "Convergence immédiate de tous les canaux (WhatsApp, Meta, Web, SMS) vers un seul flux d'onboarding."},
        {"title": "ROI & CAC optimisé", "description": "Réduction substantielle du CAC (Coût d'Acquisition Client) et hausse nette du taux de conversion."}
    ]
    add_context_slide(prs, "1. Alignement Stratégique & Positionnement", overview, points, category, 2)
    
    # Slide 3: Product Blueprint
    blueprint_intro = (
        "Le développement du MVP est régi par six principes fondamentaux garantissant la création "
        "rapide de valeur métier sans verrous techniques."
    )
    add_grid_slide(
        prs=prs,
        title="2. Product Blueprint & Principes Directeurs",
        intro_text=blueprint_intro,
        items=data["product_blueprint"]["principles"],
        category=category,
        page_num=3
    )
    
    # Slide 4: Architecture
    arch_title = "3. Architecture Technique & Socle Commun"
    arch_left_title = "Architecture Orientée Événements"
    arch_left_text = (
        f"{data['technical_architecture']['overview']}\n\n"
        f"Intégration SI : {data['technical_architecture']['integration']}"
    )
    arch_right_title = "Composants Technologiques Clés"
    # Convert string list to readable blocks
    arch_right_text = data["technical_architecture"]["core_components"].replace("1. ", "• ").replace("2. ", "\n\n• ").replace("3. ", "\n\n• ")
    add_split_card_slide(prs, arch_title, arch_left_title, arch_left_text, arch_right_title, arch_right_text, category, 4)
    
    # Slide 5-12: Capabilities (8 slides)
    page_start = 5
    for idx, cap in enumerate(data["capacities"]):
        cap_title = f"4.{idx+1} Capacité : {cap['name']}"
        
        # Client section text
        client_text = (
            f"Vision & Objectif :\n{cap['client_section']['vision']}\n\n"
            f"Valeur Métier :\n{cap['client_section']['value_added']}\n\n"
            f"Parcours Prospect :\n{cap['client_section']['use_case']}\n\n"
            f"Impact Utilisateur :\n{cap['client_section']['user_impact']}"
        )
        
        # Internal section text
        internal_text = (
            f"Données Nécessaires MVP :\n{cap['internal_team_section']['data_needed']}\n\n"
            f"Stack Technique MVP :\n{cap['internal_team_section']['tech_stack']}\n\n"
            f"APIs & Services :\n{cap['internal_team_section']['apis']}\n\n"
            f"Étapes Clés de Développement :\n{cap['internal_team_section']['key_steps']}"
        )
        
        add_split_card_slide(
            prs=prs,
            title=cap_title,
            left_title="Cadrage Client & Valeur Métier",
            left_text=client_text,
            right_title="Fiche Implémentation Interne (MVP)",
            right_text=internal_text,
            category=category,
            page_num=page_start + idx
        )
        
    # Slide 13: Roadmap
    roadmap_headers = ["Phase", "Accélérateur / Module", "Périmètre Fonctionnel", "Objectif Métier / Quick Win"]
    roadmap_data = [
        ["V1", "Digital Sales Assistant", "Smart Triage (Capacité 1) & Nurturing (Capacité 6) sur produit simple Auto", "ROI rapide, validation du canal WhatsApp"],
        ["V2", "Agent Copilot courtier", "RAG & Knowledge Hub (Capacité 4) pour le réseau d'agents physiques", "Hausse de la productivité du réseau commercial"],
        ["V3", "Conseil & Offres", "Next Best Offer (Capacité 2) & Product Advisor contextuel", "Augmentation du panier moyen de vente"],
        ["V4", "Closing & Émission", "Sentiment Analysis (Capacité 7) & Intégration Signature/Core", "Rapprocher le closing de l'automatisation 100%"],
        ["V5-V7", "Omnicanal Complet", "Intégration Open Insurance, orchestration totale des campagnes", "Écosystème prescriptif global"]
    ]
    add_table_slide(prs, "5. Roadmap Produit & Trajectoire de Maturité", roadmap_headers, roadmap_data, [1.0, 2.5, 5.0, 3.633], category, 13)
    
    # Slide 14: Business Case
    roi_headers = ["Indicateur Métier", "Performance Actuelle", "Objectif MVP (V1)", "Impact Financier Estimé"]
    roi_data = [
        ["Volume annuel estimé", "100 000 leads", "100 000 leads", "Base d'analyse"],
        ["Taux de conversion lead", "4.0%", "5.5% (+1.5 pts)", "+1 500 Ventes additionnelles / an"],
        ["Coût d'Acquisition Client", "50 €", "35 € (-30%)", "+1.5 M€ d'économies de budget marketing"],
        ["Valeur Annuelle Créée", "—", "—", "2 100 000 € / an cumulés"],
        ["Seuil de Rentabilité", "Budget MVP: 77 k€", "Point Mort : 6,5 mois", "Payback ultra-rapide sur gain de conversion"]
    ]
    add_table_slide(prs, "6. Business Case & Analyse de ROI", roi_headers, roi_data, [2.5, 2.5, 3.0, 4.133], category, 14)
    
    # Slide 15: Governance
    gov_title = "7. Modèle de Livraison & Gouvernance"
    gov_left_title = "Gouvernance & Méthode Agile"
    gov_left_text = (
        "• Comité de pilotage bimensuel (COMOP) pour arbitrer les jalons fonctionnels.\n\n"
        "• Ateliers de cadrage hebdomadaires avec les équipes sécurité, data et marketing.\n\n"
        "• Organisation Scrum en sprints de 2 semaines avec démos régulières.\n\n"
        "• Cadrage et spécifications détaillées sur les 2 premières semaines."
    )
    gov_right_title = "Conditions de Succès Clés (SanlamAllianz)"
    gov_right_text = (
        "• Validation rapide du pays pilote et de l'entité locale (Auto retail).\n\n"
        "• Mise à disposition des documents de garanties, exclusions et FAQ en format PDF.\n\n"
        "• Accès sandbox aux comptes Meta Business Manager et API WhatsApp Business.\n\n"
        "• Désignation d'un référent conformité RGPD dès le lancement."
    )
    add_split_card_slide(prs, gov_title, gov_left_title, gov_left_text, gov_right_title, gov_right_text, category, 15)
    
    # Slide 16: Closing
    create_closing_slide(prs, category)
    
    # Save
    file_path = "c:\\Projects\\Ekinox IA\\Ekinox_Digital_Sales_Accelerator_Presentation.pptx"
    prs.save(file_path)
    print(f"Sales Presentation saved to {file_path}")

def build_fraud_presentation(data):
    """Compile the PPTX deck for the Fraud Intelligence Accelerator."""
    print("Building PPTX for Ekinox Fraud Intelligence Accelerator...")
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    category = "FRAUD INTELLIGENCE ACCELERATOR"
    
    # Slide 1: Cover
    create_title_slide(
        prs=prs,
        title=data["title"],
        subtitle="Sécuriser la sinistralité et réduire le Loss Ratio via une plateforme de détection intégrée",
        category="Ekinox Insurance Platform"
    )
    
    # Slide 2: Context & Positioning
    overview = (
        "La maîtrise du loss ratio est le premier levier de rentabilité technique de SanlamAllianz.\n\n"
        "La Fraud Intelligence Platform permet de systématiser la détection des anomalies, "
        "d'optimiser le temps des enquêteurs via des alertes explicitables, et d'organiser "
        "la capitalisation continue des nouveaux modus operandi."
    )
    points = [
        {"title": "Détection précoce", "description": "Fast-track des dossiers sains et priorisation automatique des alertes suspectes selon leur montant financier."},
        {"title": "Explicabilité (Explainable AI)", "description": "Justification textuelle claire de chaque score de fraude pour accompagner la décision humaine."},
        {"title": "Boucle de capitalisation", "description": "Chaque dossier résolu enrichit automatiquement le moteur de règles pour améliorer les futures détections."}
    ]
    add_context_slide(prs, "1. Alignement Stratégique & Positionnement", overview, points, category, 2)
    
    # Slide 3: Product Blueprint
    blueprint_intro = (
        "Le développement de la plateforme anti-fraude repose sur des principes "
        "clairs de traçabilité, d'aide à la décision et d'évolution modulaire."
    )
    add_grid_slide(
        prs=prs,
        title="2. Product Blueprint & Principes Directeurs",
        intro_text=blueprint_intro,
        items=data["product_blueprint"]["principles"],
        category=category,
        page_num=3
    )
    
    # Slide 4: Architecture
    arch_title = "3. Architecture Technique & Socle Commun"
    arch_left_title = "Architecture Applicative Modulaire"
    arch_left_text = (
        f"{data['technical_architecture']['overview']}\n\n"
        f"Intégration SI : {data['technical_architecture']['integration']}"
    )
    arch_right_title = "Composants Technologiques Clés"
    arch_right_text = data["technical_architecture"]["core_components"].replace("1. ", "• ").replace("2. ", "\n\n• ").replace("3. ", "\n\n• ")
    add_split_card_slide(prs, arch_title, arch_left_title, arch_left_text, arch_right_title, arch_right_text, category, 4)
    
    # Slide 5-11: Capabilities (7 slides)
    page_start = 5
    for idx, cap in enumerate(data["capacities"]):
        cap_title = f"4.{idx+1} Capacité : {cap['name']}"
        
        client_text = (
            f"Vision & Objectif :\n{cap['client_section']['vision']}\n\n"
            f"Valeur Métier :\n{cap['client_section']['value_added']}\n\n"
            f"Cas d'Usage Enquêteur :\n{cap['client_section']['use_case']}\n\n"
            f"Impact Utilisateur :\n{cap['client_section']['user_impact']}"
        )
        
        internal_text = (
            f"Données Nécessaires MVP :\n{cap['internal_team_section']['data_needed']}\n\n"
            f"Stack Technique MVP :\n{cap['internal_team_section']['tech_stack']}\n\n"
            f"APIs & Services :\n{cap['internal_team_section']['apis']}\n\n"
            f"Étapes Clés de Développement :\n{cap['internal_team_section']['key_steps']}"
        )
        
        add_split_card_slide(
            prs=prs,
            title=cap_title,
            left_title="Cadrage Client & Valeur Métier",
            left_text=client_text,
            right_title="Fiche Implémentation Interne (MVP)",
            right_text=internal_text,
            category=category,
            page_num=page_start + idx
        )
        
    # Slide 12: Roadmap
    roadmap_headers = ["Phase", "Objectifs & Périmètre Fonctionnel", "Actions & Livrables Clés", "Valeur Métier"]
    roadmap_data = [
        ["Phase 1", "Fast-Track & Tri", "Détection des anomalies (Capacité 2) sur Auto matériel", "Fast-track immédiat des dossiers sains (OPEX)"],
        ["Phase 2", "Investigation Workspace", "Explainable AI (Capacité 5) et Timeline des dossiers", "Réduction du temps de traitement des enquêtes"],
        ["Phase 3", "Réseaux & Base", "Fraud Network Explorer (Capacité 4) & Knowledge Base (Capacité 1)", "Détection des fraudes organisées (garagistes, courtiers)"],
        ["Phase 4", "Boucle de capitalisation", "Feedback loop (Capacité 7) & réentraînement semi-automatisé", "Auto-amélioration du taux de détection avec le temps"]
    ]
    add_table_slide(prs, "5. Roadmap Produit & Boucle de Capitalisation", roadmap_headers, roadmap_data, [1.2, 3.8, 3.8, 3.333], category, 12)
    
    # Slide 13: Business Case
    roi_headers = ["Métrique ROI", "Référence Actuelle (SanlamAllianz)", "Cible MVP (V1)", "Gains Directs Annuels"]
    roi_data = [
        ["Charge Sinistre Annuelle", "100 000 000 €", "98 500 000 €", "1,5 M€ de pertes fraude évitées (-30% fuite)"],
        ["Faux Positifs de scoring", "30.0%", "<15.0%", "Productivité : qualification 2x plus rapide"],
        ["Tri Fast-Track", "Manuel", "25% de tri automatisé", "500 k€ d'OPEX économisés en gestion"],
        ["Total Gains Générés", "—", "—", "2 000 000 € / an cumulés"],
        ["Temps de retour", "Budget MVP: 60 k€", "Point Mort : 5 mois", "Rentabilité démontrée dès le 1er trimestre"]
    ]
    add_table_slide(prs, "6. Business Case & Analyse de ROI", roi_headers, roi_data, [2.5, 2.8, 2.8, 4.033], category, 13)
    
    # Slide 14: Governance
    gov_title = "7. Modèle de Livraison & Gouvernance"
    gov_left_title = "Gouvernance Projet"
    gov_left_text = (
        "• Comité de pilotage mensuel avec la direction technique et sinistres.\n\n"
        "• Ateliers bi-hebdomadaires avec les gestionnaires fraude (boucle métier).\n\n"
        "• Validation conjointe de la conformité RGPD et des critères de tri d'alertes.\n\n"
        "• Durée indicative du MVP : 8 à 10 semaines."
    )
    gov_right_title = "Conditions de Succès Clés (SanlamAllianz)"
    gov_right_text = (
        "• Extraction rapide d'un historique de sinistres qualifiés (sains et fraudes).\n\n"
        "• Accès aux bases de données des sinistres et des garages partenaires.\n\n"
        "• Disponibilité de 2 enquêteurs fraude seniors pour calibrer le workspace.\n\n"
        "• Implication active des équipes IT/Data pour les connecteurs d'extraction."
    )
    add_split_card_slide(prs, gov_title, gov_left_title, gov_left_text, gov_right_title, gov_right_text, category, 14)
    
    # Slide 15: Closing
    create_closing_slide(prs, category)
    
    # Save
    file_path = "c:\\Projects\\Ekinox IA\\Ekinox_Fraud_Intelligence_Accelerator_Presentation.pptx"
    prs.save(file_path)
    print(f"Fraud Presentation saved to {file_path}")

def main():
    json_path = "c:\\Projects\\Ekinox IA\\final_proposals.json"
    if not os.path.exists(json_path):
        print(f"Error: JSON file not found at {json_path}")
        return
        
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
        
    build_sales_presentation(data["accelerators"]["digital_sales"])
    build_fraud_presentation(data["accelerators"]["fraud_intelligence"])
    print("PowerPoint presentations compiled successfully.")

if __name__ == "__main__":
    main()
